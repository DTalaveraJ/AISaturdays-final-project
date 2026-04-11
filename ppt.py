from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(15,  23,  42)   # slide background
NAVY_CARD  = RGBColor(30,  41,  59)   # card / sidebar
TEAL       = RGBColor(20, 184, 166)   # primary accent
TEAL_DARK  = RGBColor(13, 148, 136)   # accent hover
SLATE      = RGBColor(100, 116, 139)  # muted text
WHITE      = RGBColor(255, 255, 255)
LIGHT      = RGBColor(226, 232, 240)  # body text
GOLD       = RGBColor(250, 204, 21)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

# ── Low-level helpers ─────────────────────────────────────────────────────────
def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def txbox(slide, x, y, w, h,
          text="", size=Pt(12), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def bullet_block(slide, x, y, w, h, items, size=Pt(16)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for label, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(8)
        r1 = p.add_run()
        r1.text = label
        r1.font.bold  = True
        r1.font.size  = size
        r1.font.color.rgb = TEAL
        r2 = p.add_run()
        r2.text = f"  {body}"
        r2.font.size  = size
        r2.font.color.rgb = LIGHT

def card(slide, x, y, w, h,
         title="", title_size=Pt(13), body="", body_size=Pt(12),
         accent=None):
    accent = accent or TEAL
    rect(slide, x, y, w, h, NAVY_CARD)
    rect(slide, x, y, w, Inches(0.06), accent)
    txbox(slide, x + Inches(0.18), y + Inches(0.12),
          w - Inches(0.36), Inches(0.35),
          text=title, size=title_size, bold=True, color=accent)
    txbox(slide, x + Inches(0.18), y + Inches(0.48),
          w - Inches(0.36), h - Inches(0.6),
          text=body, size=body_size, color=LIGHT, wrap=True)

# ── Content slide template ────────────────────────────────────────────────────
def slide_content(prs, title_text, build_fn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, W, Inches(1.15), NAVY_CARD)
    rect(slide, 0, Inches(1.15), W, Inches(0.05), TEAL)
    rect(slide, 0, 0, Inches(0.35), H, TEAL)
    txbox(slide,
          Inches(0.65), Inches(0.22), Inches(11.5), Inches(0.75),
          text=title_text, size=Pt(28), bold=True, color=WHITE)
    build_fn(slide)
    return slide

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)

    # Left teal sidebar
    rect(slide, 0, 0, Inches(0.55), H, TEAL)

    # Decorative top-right block
    rect(slide, Inches(10.5), 0, Inches(2.83), Inches(1.8), TEAL_DARK)

    # Bottom-right decorative card
    rect(slide, Inches(9.5), Inches(5.8), Inches(3.83), Inches(1.7), NAVY_CARD)

    # Tag line
    txbox(slide,
          Inches(1.1), Inches(1.6), Inches(8), Inches(0.5),
          text="INTELIGENCIA ARTIFICIAL · CONSUMO EFICIENTE",
          size=Pt(11), bold=True, color=TEAL)

    # Main title
    box = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(8.8), Inches(2.2))
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Tu AIbuela\nde Confianza"
    r.font.size  = Pt(54)
    r.font.bold  = True
    r.font.color.rgb = WHITE

    # Subtitle in teal
    txbox(slide,
          Inches(1.1), Inches(4.45), Inches(9), Inches(0.55),
          text="Proyecto ShopPath IA · Optimización de Consumo Inteligente",
          size=Pt(22), bold=True, color=TEAL)

    # Divider
    rect(slide, Inches(1.1), Inches(5.15), Inches(4.5), Inches(0.04), SLATE)

    # Date / descriptor
    txbox(slide,
          Inches(1.1), Inches(5.3), Inches(8), Inches(0.45),
          text="Eficiencia algorítmica aplicada al ahorro doméstico  ·  21 de Marzo, 2026",
          size=Pt(13), color=SLATE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MISIÓN
# ═════════════════════════════════════════════════════════════════════════════
def build_mision(slide):
    items = [
        ("El Problema:",
         "Disparidad de precios y coste de oportunidad: el consumidor pierde dinero y tiempo por falta de información."),
        ("Nuestra Misión:",
         "Democratizar los datos de consumo para generar ahorro real y medible para cualquier hogar."),
        ("El Objetivo:",
         "Maximizar el poder adquisitivo equilibrando precio, distancia y número de paradas."),
    ]
    bullet_block(slide, Inches(0.65), Inches(1.4), Inches(8.4), Inches(5.5), items, size=Pt(17))

    card(slide,
         Inches(9.3), Inches(1.55), Inches(3.5), Inches(1.5),
         title="Ahorro potencial", body="15 – 25 %\npor hogar / mes",
         title_size=Pt(12), body_size=Pt(20))

    card(slide,
         Inches(9.3), Inches(3.25), Inches(3.5), Inches(1.5),
         title="Nombre del proyecto", body="Tu AIbuela\nde Confianza",
         title_size=Pt(12), body_size=Pt(18))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARQUITECTURA
# ═════════════════════════════════════════════════════════════════════════════
def build_arquitectura(slide):
    cards = [
        ("Ingesta de datos",  "Python + Playwright\npara catálogos digitales"),
        ("IA · NLP",          "Llama 3 para normalizar\nlistas del usuario"),
        ("Geo & Routing",     "OpenStreetMap · OSRM\nrutas de precisión"),
        ("Cloud & Backend",   "Supabase + Render\ninfraestructura escalable"),
    ]
    col_w = Inches(2.9)
    gap   = Inches(0.28)
    for i, (t, b) in enumerate(cards):
        x = Inches(0.65) + i * (col_w + gap)
        card(slide, x, Inches(1.45), col_w, Inches(3.8),
             title=t, body=b, title_size=Pt(13), body_size=Pt(15))

    txbox(slide,
          Inches(0.65), Inches(5.55), Inches(12.05), Inches(0.5),
          text="100 % Open Source  ·  sin dependencias propietarias  ·  coste operativo mínimo",
          size=Pt(12), color=SLATE, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ALGORITMO
# ═════════════════════════════════════════════════════════════════════════════
def build_algoritmo(slide):
    steps = [
        ("01  Coste",          "Selección del mínimo precio unitario por SKU en el catálogo combinado."),
        ("02  Desplazamiento", "Minimización de distancia total y huella de carbono asociada al trayecto."),
        ("03  K-Stops",        "Restricción inteligente del nº de establecimientos para rentabilizar la ruta."),
        ("Regla de Oro",       "El ahorro en precio debe superar siempre el coste del trayecto."),
    ]
    for i, (label, body) in enumerate(steps):
        yy = Inches(1.42) + i * Inches(1.32)
        rect(slide, Inches(0.65), yy, Inches(2.5), Inches(1.05), NAVY_CARD)
        rect(slide, Inches(0.65), yy, Inches(0.08), Inches(1.05), TEAL)
        txbox(slide, Inches(0.85), yy + Inches(0.12),
              Inches(2.2), Inches(0.4),
              text=label, size=Pt(13), bold=True, color=TEAL)
        txbox(slide, Inches(3.35), yy + Inches(0.1),
              Inches(9.35), Inches(0.9),
              text=body, size=Pt(15), color=LIGHT)
        if i < 3:
            rect(slide, Inches(0.65), yy + Inches(1.08),
                 Inches(12.05), Inches(0.02), NAVY_CARD)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROADMAP
# ═════════════════════════════════════════════════════════════════════════════
def build_roadmap(slide):
    phases = [
        ("Semanas 1–2",  "MVP de extracción\nde datos locales"),
        ("Semana 3",     "Integración de IA\npara listas inteligentes"),
        ("Semana 4",     "Motor de rutas\nvisuales en vivo"),
        ("Resultado",    "15–25 % de ahorro\nmensual garantizado"),
    ]
    col_w = Inches(2.8)
    gap   = Inches(0.3)
    for i, (phase, desc) in enumerate(phases):
        x      = Inches(0.65) + i * (col_w + gap)
        accent = GOLD if i == 3 else TEAL
        rect(slide, x, Inches(1.45), col_w, Inches(2.9), NAVY_CARD)
        rect(slide, x, Inches(1.45), col_w, Inches(0.06), accent)
        txbox(slide, x + Inches(0.15), Inches(1.57),
              col_w - Inches(0.3), Inches(0.4),
              text=phase, size=Pt(12), bold=True, color=accent)
        txbox(slide, x + Inches(0.15), Inches(2.05),
              col_w - Inches(0.3), Inches(2.1),
              text=desc, size=Pt(15), color=LIGHT)
        if i < 3:
            ax = x + col_w + Inches(0.05)
            txbox(slide, ax, Inches(2.7), gap, Inches(0.4),
                  text="→", size=Pt(20), bold=True, color=TEAL,
                  align=PP_ALIGN.CENTER)

    txbox(slide,
          Inches(0.65), Inches(4.65), Inches(12.05), Inches(2.4),
          text=(
              "El roadmap está diseñado para validar hipótesis rápido (2 semanas) y escalar "
              "progresivamente. Cada fase entrega valor independiente, minimizando el riesgo "
              "de inversión y permitiendo pivotar con datos reales."
          ),
          size=Pt(13), color=SLATE, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
def create_presentation():
    prs = new_prs()

    slide_cover(prs)
    slide_content(prs, "Misión y Visión del Proyecto",         build_mision)
    slide_content(prs, "Ecosistema Tecnológico (Open Source)", build_arquitectura)
    slide_content(prs, "Motor de Optimización Multiobjetivo",  build_algoritmo)
    slide_content(prs, "Impacto y Roadmap de Implementación",  build_roadmap)

    prs.save("Proyecto_ShopPath_IA.pptx")
    print("✓ Presentación guardada: Proyecto_ShopPath_IA.pptx")

create_presentation()
